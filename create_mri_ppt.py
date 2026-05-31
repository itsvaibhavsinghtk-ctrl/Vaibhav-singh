from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x0D, 0x2B, 0x55)
MID_BLUE  = RGBColor(0x1A, 0x5F, 0x9E)
ACCENT    = RGBColor(0x00, 0xC2, 0xCB)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, title, bullets, l, t, w, h, bg_color=LIGHT_GRAY, title_color=DARK_BLUE):
    add_rect(slide, l, t, w, h, bg_color)
    add_text_box(slide, title, l+0.1, t+0.08, w-0.2, 0.4, 13, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(Inches(l+0.1), Inches(t+0.5), Inches(w-0.2), Inches(h-0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(10.5)
        run.font.color.rgb = DARK_TEXT

# ─── SLIDE 1: Title Slide ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 2.8, 13.33, 0.04, ACCENT)
add_rect(slide, 0, 5.0, 13.33, 0.04, MID_BLUE)
add_text_box(slide, "MRI SEQUENCES", 0.5, 0.8, 12.3, 1.2, 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "A Comprehensive Guide for Clinical Radiology", 0.5, 2.0, 12.3, 0.7, 20, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, "Understanding T1, T2, FLAIR, DWI, GRE, STIR & More", 0.5, 3.0, 12.3, 0.6, 15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, "Radiology / Medical Imaging Lecture Series", 1.0, 5.2, 11.3, 0.5, 13, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, "2026", 1.0, 5.8, 11.3, 0.5, 13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ─── SLIDE 2: Table of Contents ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
add_rect(slide, 0, 1.1, 0.08, 6.4, ACCENT)
add_text_box(slide, "TABLE OF CONTENTS", 0.3, 0.25, 12.0, 0.7, 28, bold=True, color=WHITE)

topics = [
    ("01", "MRI Physics Basics", "How MRI signal is generated"),
    ("02", "T1-Weighted Imaging", "Short TR, Short TE — anatomy & post-contrast"),
    ("03", "T2-Weighted Imaging", "Long TR, Long TE — pathology detection"),
    ("04", "FLAIR Sequence", "Fluid Attenuated Inversion Recovery"),
    ("05", "DWI & ADC", "Diffusion Weighted Imaging"),
    ("06", "GRE / T2* Sequences", "Gradient Echo & susceptibility"),
    ("07", "STIR Sequence", "Short Tau Inversion Recovery — fat suppression"),
    ("08", "Special Sequences", "MRA, MRS, Perfusion, fMRI & more"),
    ("09", "Clinical Applications", "Choosing the right sequence"),
]

cols = [0.3, 4.55, 8.8]
for i, (num, title, desc) in enumerate(topics):
    col = i // 3
    row = i % 3
    x = cols[col]
    y = 1.3 + row * 1.9
    add_rect(slide, x, y, 4.0, 1.6, WHITE)
    add_rect(slide, x, y, 0.55, 1.6, MID_BLUE)
    add_text_box(slide, num, x+0.05, y+0.4, 0.5, 0.7, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x+0.65, y+0.1, 3.2, 0.5, 12, bold=True, color=DARK_BLUE)
    add_text_box(slide, desc, x+0.65, y+0.6, 3.2, 0.7, 9.5, color=DARK_TEXT, italic=True)


# ─── SLIDE 3: MRI Physics Basics ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_rect(slide, 0, 1.2, 13.33, 0.06, ACCENT)
add_text_box(slide, "01  MRI PHYSICS BASICS", 0.3, 0.25, 12.0, 0.8, 26, bold=True, color=WHITE)

add_bullet_box(slide, "How MRI Signal is Generated", [
    "Hydrogen protons (¹H) align with external magnetic field (B₀)",
    "Radiofrequency (RF) pulse tips protons out of alignment",
    "Protons precess at Larmor frequency: ω = γ × B₀",
    "Signal detected as protons relax back to equilibrium",
    "Two relaxation processes: T1 (longitudinal) & T2 (transverse)",
], 0.3, 1.4, 6.0, 3.5, LIGHT_GRAY, DARK_BLUE)

add_bullet_box(slide, "Key Parameters", [
    "TR (Repetition Time) — time between RF pulses (ms)",
    "TE (Echo Time) — time between RF pulse and signal reading",
    "TI (Inversion Time) — used in inversion recovery sequences",
    "Flip Angle — degree of magnetization tipping",
    "Field Strength: 1.5T (standard) vs 3T (high-res)",
], 6.5, 1.4, 6.5, 3.5, LIGHT_GRAY, DARK_BLUE)

add_rect(slide, 0.3, 5.1, 12.7, 1.8, RGBColor(0xE8, 0xF4, 0xFF))
add_text_box(slide, "⚡ Key Concept", 0.5, 5.15, 3.0, 0.4, 12, bold=True, color=MID_BLUE)
add_text_box(slide, "T1 relaxation = longitudinal recovery (spin-lattice). T2 relaxation = transverse decay (spin-spin). "
             "Sequence contrast depends on TR and TE choices relative to tissue T1/T2 values.",
             0.5, 5.55, 12.3, 1.2, 11, color=DARK_TEXT)


# ─── SLIDE 4: T1-Weighted Imaging ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x8B, 0x45, 0x13))
add_rect(slide, 0, 1.2, 13.33, 0.06, ORANGE)
add_text_box(slide, "02  T1-WEIGHTED IMAGING", 0.3, 0.25, 12.0, 0.8, 26, bold=True, color=WHITE)

add_bullet_box(slide, "Parameters", [
    "Short TR: 400–600 ms",
    "Short TE: 10–30 ms",
    "Emphasizes T1 differences between tissues",
], 0.3, 1.4, 4.0, 2.8, RGBColor(0xFF, 0xF3, 0xE8), RGBColor(0x8B, 0x45, 0x13))

add_bullet_box(slide, "Tissue Appearance", [
    "FAT → Bright (high signal) ✨",
    "FLUID (CSF, urine) → Dark (low signal)",
    "MUSCLE → Intermediate",
    "CORTICAL BONE → Dark",
    "GREY MATTER → Slightly darker than white matter",
    "WHITE MATTER → Brighter than grey matter",
], 4.5, 1.4, 4.3, 2.8, RGBColor(0xFF, 0xF3, 0xE8), RGBColor(0x8B, 0x45, 0x13))

add_bullet_box(slide, "Clinical Uses", [
    "Anatomy — excellent soft tissue detail",
    "Post-gadolinium contrast enhancement",
    "Fat detection (lipoma, dermoid)",
    "Hemorrhage (subacute — methemoglobin bright)",
    "Liver lesions characterization",
    "Pituitary gland imaging",
    "Bone marrow assessment",
], 9.0, 1.4, 4.0, 2.8, RGBColor(0xFF, 0xF3, 0xE8), RGBColor(0x8B, 0x45, 0x13))

add_rect(slide, 0.3, 4.4, 12.7, 2.7, RGBColor(0xFF, 0xED, 0xD5))
add_text_box(slide, "🔶 T1 BRIGHT (Short T1)", 0.5, 4.45, 6.0, 0.45, 12, bold=True, color=ORANGE)
add_text_box(slide, "Fat  |  Subacute blood (Met-Hb)  |  Gadolinium contrast  |  Proteinaceous fluid  |  Melanin  |  Calcium (sometimes)  |  Lipid-containing lesions",
             0.5, 4.9, 12.3, 0.6, 10.5, color=DARK_TEXT)
add_text_box(slide, "🔷 T1 DARK (Long T1)", 0.5, 5.6, 6.0, 0.45, 12, bold=True, color=MID_BLUE)
add_text_box(slide, "Free fluid (CSF, edema, cysts)  |  Acute blood  |  Air  |  Cortical bone  |  Most pathology (tumors, infarcts)",
             0.5, 6.05, 12.3, 0.6, 10.5, color=DARK_TEXT)


# ─── SLIDE 5: T2-Weighted Imaging ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, MID_BLUE)
add_rect(slide, 0, 1.2, 13.33, 0.06, ACCENT)
add_text_box(slide, "03  T2-WEIGHTED IMAGING", 0.3, 0.25, 12.0, 0.8, 26, bold=True, color=WHITE)

add_bullet_box(slide, "Parameters", [
    "Long TR: 2000–4000 ms",
    "Long TE: 80–120 ms",
    "Emphasizes T2 differences — fluid appears bright",
    "Called the 'pathology sequence'",
], 0.3, 1.4, 4.1, 3.0, RGBColor(0xE8, 0xF4, 0xFF), MID_BLUE)

add_bullet_box(slide, "Tissue Appearance", [
    "FLUID (CSF, edema) → Bright ✨",
    "FAT → Intermediate-bright",
    "MUSCLE → Dark",
    "CORTICAL BONE → Dark",
    "GREY MATTER → Brighter than white matter",
    "WHITE MATTER → Darker than grey matter",
    "PATHOLOGY → Mostly bright (edema, tumors)",
], 4.6, 1.4, 4.3, 3.0, RGBColor(0xE8, 0xF4, 0xFF), MID_BLUE)

add_bullet_box(slide, "Clinical Uses", [
    "Edema detection (brain, spine, joints)",
    "Joint fluid & cartilage evaluation",
    "Spinal cord pathology",
    "Liver & kidney lesions",
    "Pelvic organs (uterus, prostate)",
    "Breast MRI",
    "Inflammation & infection",
    "CSF flow studies",
], 9.1, 1.4, 4.0, 3.0, RGBColor(0xE8, 0xF4, 0xFF), MID_BLUE)

add_rect(slide, 0.3, 4.6, 12.7, 2.6, RGBColor(0xD6, 0xEA, 0xFF))
add_text_box(slide, "T1 vs T2 — Quick Comparison", 0.5, 4.65, 12.0, 0.4, 12, bold=True, color=DARK_BLUE)
headers = ["Tissue", "T1 Signal", "T2 Signal"]
data = [["Fat","Bright","Bright"],["Free Fluid","Dark","Bright"],["Muscle","Intermediate","Dark"],["Air/Bone","Dark","Dark"],["Edema","Dark","Bright"],["Subacute Blood","Bright","Bright"]]
col_x = [0.4, 3.5, 7.0]
for ci, h in enumerate(headers):
    add_text_box(slide, h, col_x[ci], 5.1, 3.0, 0.35, 10, bold=True, color=MID_BLUE)
for ri, row in enumerate(data):
    y = 5.5 + ri * 0.28
    for ci, cell in enumerate(row):
        clr = GREEN if cell == "Bright" else (RGBColor(0xCC,0x00,0x00) if cell == "Dark" else DARK_TEXT)
        add_text_box(slide, cell, col_x[ci], y, 3.0, 0.28, 9.5, color=clr)


# ─── SLIDE 6: FLAIR ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x4A, 0x00, 0x8C))
add_rect(slide, 0, 1.2, 13.33, 0.06, RGBColor(0xB0, 0x5C, 0xFF))
add_text_box(slide, "04  FLAIR — Fluid Attenuated Inversion Recovery", 0.3, 0.25, 12.5, 0.8, 23, bold=True, color=WHITE)

add_bullet_box(slide, "How FLAIR Works", [
    "Inversion Recovery sequence with long TI (~2000 ms at 1.5T)",
    "TI chosen to null free water/CSF signal",
    "CSF appears BLACK (suppressed)",
    "Periventricular & cortical lesions become visible",
    "Lesions near CSF spaces that are hidden on T2 become conspicuous",
], 0.3, 1.4, 6.2, 3.2, RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0x4A, 0x00, 0x8C))

add_bullet_box(slide, "Key Clinical Applications", [
    "Multiple Sclerosis (MS) plaques — periventricular",
    "Subarachnoid hemorrhage (SAH) — blood bright in CSF spaces",
    "Cortical contusions & traumatic brain injury",
    "Encephalitis (herpes, autoimmune)",
    "Leptomeningeal disease / carcinomatosis",
    "Early ischemic stroke changes",
    "Cortical dysplasia",
], 6.7, 1.4, 6.3, 3.2, RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0x4A, 0x00, 0x8C))

add_rect(slide, 0.3, 4.8, 12.7, 2.4, RGBColor(0xF0, 0xE6, 0xFF))
add_text_box(slide, "⚡ FLAIR vs T2 — When to use which?", 0.5, 4.85, 12.0, 0.4, 12, bold=True, color=RGBColor(0x4A, 0x00, 0x8C))
add_text_box(slide, "T2: Best overall pathology detection — fluid AND lesions both bright\n"
             "FLAIR: Best for PERIVENTRICULAR and CORTICAL lesions — nulls CSF so adjacent lesions stand out\n"
             "Rule: Always review T2 and FLAIR together for brain pathology (especially MS, SAH, encephalitis)",
             0.5, 5.3, 12.3, 1.7, 10.5, color=DARK_TEXT)


# ─── SLIDE 7: DWI & ADC ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x00, 0x6B, 0x3C))
add_rect(slide, 0, 1.2, 13.33, 0.06, GREEN)
add_text_box(slide, "05  DWI & ADC — Diffusion Weighted Imaging", 0.3, 0.25, 12.5, 0.8, 24, bold=True, color=WHITE)

add_bullet_box(slide, "Principle of DWI", [
    "Measures Brownian motion of water molecules",
    "Applies gradient pulses to sensitize motion",
    "b-value controls diffusion weighting (b=0, 500, 1000 s/mm²)",
    "RESTRICTED diffusion → Bright on DWI",
    "FREE diffusion → Dark on DWI",
    "T2 shine-through can mimic restriction — always check ADC!",
], 0.3, 1.4, 6.2, 3.2, RGBColor(0xE8, 0xFF, 0xF0), RGBColor(0x00, 0x6B, 0x3C))

add_bullet_box(slide, "ADC Map Interpretation", [
    "ADC = Apparent Diffusion Coefficient map",
    "Mathematically derived from DWI",
    "TRUE restriction: DWI bright + ADC dark (low ADC)",
    "T2 shine-through: DWI bright + ADC bright",
    "Low ADC value < 0.8 × 10⁻³ mm²/s = significant restriction",
    "High ADC = free diffusion / vasogenic edema",
], 6.7, 1.4, 6.3, 3.2, RGBColor(0xE8, 0xFF, 0xF0), RGBColor(0x00, 0x6B, 0x3C))

add_rect(slide, 0.3, 4.8, 12.7, 2.4, RGBColor(0xE0, 0xFF, 0xEE))
add_text_box(slide, "🏥 Clinical Applications of DWI", 0.5, 4.85, 12.0, 0.4, 12, bold=True, color=RGBColor(0x00, 0x6B, 0x3C))
apps = [
    "Acute ischemic stroke (cytotoxic edema — bright within minutes)",
    "Abscess (restricted pus — bright center, helps vs necrotic tumor)",
    "Epidermoid cyst vs arachnoid cyst",
    "Diffuse axonal injury (DAI)",
    "Creutzfeldt-Jakob disease (CJD) — cortical ribboning",
    "Lymphoma (hypercellular — restricted)",
    "Cholesteatoma",
    "Prostate & rectal cancer staging",
]
for i in range(4):
    add_text_box(slide, "• " + apps[i], 0.5, 5.3 + i*0.25, 6.0, 0.28, 9.5, color=DARK_TEXT)
for i in range(4):
    add_text_box(slide, "• " + apps[i+4], 6.8, 5.3 + i*0.25, 6.0, 0.28, 9.5, color=DARK_TEXT)


# ─── SLIDE 8: GRE / T2* ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x7B, 0x1F, 0x00))
add_rect(slide, 0, 1.2, 13.33, 0.06, ORANGE)
add_text_box(slide, "06  GRE / T2* — Gradient Echo & Susceptibility", 0.3, 0.25, 12.5, 0.8, 24, bold=True, color=WHITE)

add_bullet_box(slide, "GRE Principle", [
    "Uses gradient reversal instead of 180° RF refocusing pulse",
    "Does NOT cancel field inhomogeneity effects",
    "Susceptible to T2* (T2-star) decay",
    "Much faster than spin echo sequences",
    "Flip angle < 90° (typically 15–60°)",
    "Small flip angle → more T2* weighting",
], 0.3, 1.4, 4.1, 3.2, RGBColor(0xFF, 0xEE, 0xE0), RGBColor(0x7B, 0x1F, 0x00))

add_bullet_box(slide, "Susceptibility Artifacts", [
    "Paramagnetic substances cause local field distortion",
    "Blood products (deoxyhemoglobin, hemosiderin) → DARK",
    "Air-tissue interfaces → signal loss / blooming",
    "Metal implants → severe distortion",
    "'Blooming artifact' — lesions appear larger than actual",
    "SWI (Susceptibility Weighted Imaging) = enhanced T2*",
], 4.5, 1.4, 4.4, 3.2, RGBColor(0xFF, 0xEE, 0xE0), RGBColor(0x7B, 0x1F, 0x00))

add_bullet_box(slide, "Clinical Uses", [
    "Microbleeds (hypertensive, CAA, DAI)",
    "Cavernous malformations",
    "Hemorrhagic transformation of infarct",
    "Iron deposition (basal ganglia disorders)",
    "Calcifications (T2* dark)",
    "Venous thrombosis (SWI)",
    "Tumor blood products",
    "Functional MRI (fMRI) — BOLD technique",
], 9.1, 1.4, 4.0, 3.2, RGBColor(0xFF, 0xEE, 0xE0), RGBColor(0x7B, 0x1F, 0x00))

add_rect(slide, 0.3, 4.85, 12.7, 2.3, RGBColor(0xFF, 0xF0, 0xE8))
add_text_box(slide, "Hemorrhage on MRI — Signal Evolution Over Time", 0.5, 4.9, 12.0, 0.4, 12, bold=True, color=RGBColor(0x7B, 0x1F, 0x00))
stages = [("Hyperacute\n(<24h)","Oxy-Hb","Dark","Bright","Iso/Dark"),
          ("Acute\n(1–3d)","Deoxy-Hb","Dark","Dark","Dark"),
          ("Early Subacute\n(3–7d)","Met-Hb (intra)","Bright","Dark","Dark"),
          ("Late Subacute\n(1–2wk)","Met-Hb (extra)","Bright","Bright","Dark"),
          ("Chronic\n(>2wk)","Hemosiderin","Dark","Dark","Very Dark")]
hx = [0.4, 2.9, 5.2, 7.5, 9.8, 11.8]
hlabels = ["Stage","Content","T1","T2","GRE"]
for ci, h in enumerate(hlabels):
    add_text_box(slide, h, hx[ci], 5.35, 2.3, 0.3, 9, bold=True, color=RGBColor(0x7B, 0x1F, 0x00))
for ri, row in enumerate(stages):
    y = 5.65 + ri*0.3
    for ci, cell in enumerate(row):
        add_text_box(slide, cell, hx[ci], y, 2.3, 0.3, 8.5, color=DARK_TEXT)


# ─── SLIDE 9: STIR ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x00, 0x5C, 0x6B))
add_rect(slide, 0, 1.2, 13.33, 0.06, ACCENT)
add_text_box(slide, "07  STIR — Short Tau Inversion Recovery", 0.3, 0.25, 12.5, 0.8, 24, bold=True, color=WHITE)

add_bullet_box(slide, "How STIR Works", [
    "Inversion Recovery sequence with SHORT TI",
    "TI chosen to null FAT signal (~150ms at 1.5T, ~210ms at 3T)",
    "Fat appears BLACK (suppressed)",
    "Fluid/edema remains BRIGHT",
    "Additive effect: T1 + T2 weighting → very sensitive to water",
    "Works at any field strength (unlike spectral fat sat)",
], 0.3, 1.4, 6.2, 3.2, RGBColor(0xE0, 0xF7, 0xF9), RGBColor(0x00, 0x5C, 0x6B))

add_bullet_box(slide, "Clinical Applications", [
    "Bone marrow edema (fractures, osteomyelitis, metastases)",
    "Soft tissue edema & inflammation",
    "Muscle injuries & tears",
    "Brachial plexus imaging",
    "Spinal cord & nerve root pathology",
    "Lymph node evaluation",
    "Body MRI where spectral fat sat unreliable",
    "Pediatric imaging",
], 6.7, 1.4, 6.3, 3.2, RGBColor(0xE0, 0xF7, 0xF9), RGBColor(0x00, 0x5C, 0x6B))

add_rect(slide, 0.3, 4.8, 12.7, 2.4, RGBColor(0xD0, 0xF2, 0xF5))
add_text_box(slide, "STIR vs Spectral Fat Saturation (Fat Sat) — Key Differences", 0.5, 4.85, 12.0, 0.4, 12, bold=True, color=RGBColor(0x00, 0x5C, 0x6B))
comp = [
    ("Feature", "STIR", "Spectral Fat Sat"),
    ("Mechanism", "Inversion Recovery (TI)", "Chemical shift saturation pulse"),
    ("Field strength", "Works at all field strengths", "Best at high field (1.5T, 3T)"),
    ("Uniformity", "Uniform fat suppression", "Inhomogeneous at air-tissue interfaces"),
    ("SNR", "Lower SNR", "Higher SNR"),
    ("Contrast injection", "Cannot be used (nulls Gd)", "Safe with contrast"),
]
cx = [0.5, 3.5, 8.3]
for ri, row in enumerate(comp):
    y = 5.3 + ri * 0.28
    for ci, cell in enumerate(row):
        bold = ri == 0
        add_text_box(slide, cell, cx[ci], y, 4.5, 0.28, 9 if not bold else 9.5, bold=bold,
                     color=RGBColor(0x00, 0x5C, 0x6B) if bold else DARK_TEXT)


# ─── SLIDE 10: Special Sequences ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, RGBColor(0x1A, 0x1A, 0x5E))
add_rect(slide, 0, 1.2, 13.33, 0.06, RGBColor(0x77, 0x77, 0xFF))
add_text_box(slide, "08  SPECIAL SEQUENCES", 0.3, 0.25, 12.5, 0.8, 26, bold=True, color=WHITE)

specials = [
    ("MRA\nMR Angiography", [
        "TOF (Time of Flight) — flow-related enhancement",
        "Phase Contrast — velocity encoding",
        "CE-MRA — contrast enhanced",
        "Evaluates arteries & veins non-invasively",
    ], 0.3, 1.4, 3.0, 3.2, RGBColor(0xEA, 0xEA, 0xFF)),
    ("MRS\nMR Spectroscopy", [
        "Measures brain metabolite peaks",
        "NAA (neurons), Cho (membranes),",
        "Cr (reference), Lac (anaerobic), Lip",
        "Tumor vs radiation necrosis",
    ], 3.5, 1.4, 3.0, 3.2, RGBColor(0xEA, 0xEA, 0xFF)),
    ("Perfusion MRI", [
        "DSC (Dynamic Susceptibility Contrast)",
        "ASL (Arterial Spin Labeling) — no contrast",
        "Maps CBF, CBV, MTT, TTP",
        "Stroke, tumor grading, seizures",
    ], 6.7, 1.4, 3.0, 3.2, RGBColor(0xEA, 0xEA, 0xFF)),
    ("fMRI\nFunctional MRI", [
        "BOLD technique (Blood Oxygen Level Dependent)",
        "Maps brain activation patterns",
        "Pre-surgical planning",
        "Eloquent cortex mapping",
    ], 9.9, 1.4, 3.1, 3.2, RGBColor(0xEA, 0xEA, 0xFF)),
]
for title, bullets, l, t, w, h, bg in specials:
    add_bullet_box(slide, title, bullets, l, t, w, h, bg, RGBColor(0x1A, 0x1A, 0x5E))

# Bottom row
specials2 = [
    ("MRCP", ["MR Cholangiopancreatography","Heavy T2 — fluid bright","Bile ducts, pancreatic duct","No contrast, no radiation"], 0.3, 4.8, 3.0, 2.4, RGBColor(0xF0, 0xF0, 0xFF)),
    ("DTI\nDiffusion Tensor", ["Tracks white matter tracts","Fractional Anisotropy (FA) maps","Pre-surgical tractography","TBI assessment"], 3.5, 4.8, 3.0, 2.4, RGBColor(0xF0, 0xF0, 0xFF)),
    ("Dixon / Fat-Water", ["Separates fat & water images","In-phase / Out-of-phase","PDFF (fat quantification)","Liver steatosis grading"], 6.7, 4.8, 3.0, 2.4, RGBColor(0xF0, 0xF0, 0xFF)),
    ("Cine MRI / EPI", ["Cardiac gated sequences","Real-time cardiac function","EPI — fast single-shot","fMRI & perfusion base"], 9.9, 4.8, 3.1, 2.4, RGBColor(0xF0, 0xF0, 0xFF)),
]
for title, bullets, l, t, w, h, bg in specials2:
    add_bullet_box(slide, title, bullets, l, t, w, h, bg, RGBColor(0x1A, 0x1A, 0x5E))


# ─── SLIDE 11: Clinical Applications ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
add_rect(slide, 0, 1.2, 13.33, 0.06, ACCENT)
add_text_box(slide, "09  CLINICAL APPLICATIONS — Choosing the Right Sequence", 0.3, 0.25, 12.5, 0.8, 22, bold=True, color=WHITE)

cases = [
    ("Acute Stroke", "DWI + ADC (first line)\nFLAIR (after 6h)\nMRA for vessel occlusion\nT2* / SWI for hemorrhage"),
    ("Brain Tumor", "T1 + T1 post-Gad (enhancement)\nT2 / FLAIR (extent, edema)\nDWI (cellularity)\nMRS (metabolites)\nPerfusion (grading)"),
    ("Multiple Sclerosis", "FLAIR (periventricular plaques)\nT2 (overall lesion load)\nT1 (black holes = axonal loss)\nSpinal T2 (cord lesions)"),
    ("Spine Pathology", "T1 (anatomy, bone marrow)\nT2 (disc, cord, fluid)\nSTIR (edema, metastases)\nDWI (cord, vertebral lesion)"),
    ("Knee / Joint", "PD Fat Sat / T2 Fat Sat (menisci, cartilage)\nT1 (anatomy)\nGRE (cartilage volume)\nSTIR (bone bruise, edema)"),
    ("Liver Lesion", "T1 in/out phase (fat, iron)\nT2 (cyst, hemangioma)\nDWI (malignancy)\nT1 dynamic post-Gad\nDixon fat quantification"),
]
cols_x = [0.3, 4.55, 8.8]
for i, (title, content) in enumerate(cases):
    col = i % 3
    row = i // 3
    x = cols_x[col]
    y = 1.4 + row * 2.8
    add_rect(slide, x, y, 4.0, 2.6, LIGHT_GRAY)
    add_rect(slide, x, y, 4.0, 0.45, MID_BLUE)
    add_text_box(slide, title, x+0.1, y+0.05, 3.8, 0.38, 11, bold=True, color=WHITE)
    add_text_box(slide, content, x+0.1, y+0.52, 3.8, 2.0, 9.5, color=DARK_TEXT)


# ─── SLIDE 12: Quick Reference Summary ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
add_rect(slide, 0, 7.44, 13.33, 0.06, ACCENT)
add_text_box(slide, "QUICK REFERENCE — MRI Sequences Summary", 0.3, 0.15, 12.5, 0.8, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Sequence", "TR", "TE", "Fat", "Fluid", "Key Use", True),
    ("T1", "Short", "Short", "Bright", "Dark", "Anatomy, Post-Gad, Fat", False),
    ("T2", "Long", "Long", "Bright", "Bright", "Pathology, Edema, Fluid", False),
    ("FLAIR", "Long", "Long", "Bright", "Dark (nulled)", "Periventricular lesions, SAH", False),
    ("DWI", "—", "Short", "—", "—", "Acute stroke, Abscess, Tumor", False),
    ("GRE/T2*", "Short", "Short", "Bright", "Dark", "Hemorrhage, Iron, Calcification", False),
    ("STIR", "Long", "Long", "Dark (nulled)", "Bright", "Bone marrow edema, Soft tissue", False),
    ("SWI", "Long", "Long", "—", "—", "Microbleeds, Veins, Calcium", False),
    ("MRCP", "Very Long", "Very Long", "Dark", "Very Bright", "Bile/Pancreatic ducts", False),
    ("MRS", "—", "—", "—", "—", "Tumor metabolites, Necrosis", False),
]

col_x = [0.2, 2.1, 3.3, 4.5, 5.8, 7.2, 9.5]
col_w = [1.8, 1.1, 1.1, 1.2, 1.3, 2.2, 3.5]
for ri, row_data in enumerate(rows):
    y = 1.05 + ri * 0.62
    is_header = row_data[-1]
    bg = MID_BLUE if is_header else (RGBColor(0x0F, 0x35, 0x60) if ri % 2 == 0 else RGBColor(0x13, 0x40, 0x72))
    add_rect(slide, 0.2, y, 12.9, 0.58, bg)
    for ci, (cell, cx, cw) in enumerate(zip(row_data[:-1], col_x, col_w)):
        add_text_box(slide, str(cell), cx+0.05, y+0.1, cw-0.1, 0.42,
                     10 if not is_header else 10.5,
                     bold=is_header,
                     color=ACCENT if is_header else (GREEN if cell == "Bright" else
                           (RGBColor(0xFF,0x88,0x88) if cell == "Dark" else
                            (ORANGE if cell == "Dark (nulled)" else WHITE))))

# Save
prs.save("/projects/sandbox/MRI_Sequences_Lecture.pptx")
print("✅ MRI Sequences PPT saved to /projects/sandbox/MRI_Sequences_Lecture.pptx")
